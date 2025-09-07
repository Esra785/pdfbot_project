# pdfbot_7.py
# 2. Hafta - 3. Gün 
# CSV -> PowerPoint slaytları

import pandas as pd
from pptx import Presentation

INPUT_CSV = "emsal_focused_table.csv"
OUTPUT_PPTX = "emsal_presentation.pptx"

df = pd.read_csv(INPUT_CSV).fillna({"Kategori": "Bilinmiyor", "Odaklı Özet": "Özet yok"})

prs = Presentation()

# Başlık slaytı
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
if title_slide.shapes.title:
    title_slide.shapes.title.text = "Emsal Karar Özeti"
if len(title_slide.placeholders) > 1 and title_slide.placeholders[1]:
    title_slide.placeholders[1].text = "LLM destekli analiz (otomatik)" # type: ignore

# İçerik slaytları
for _, row in df.iterrows():
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    page_info = row.get("Sayfa", "Genel")
    kategori = row.get("Kategori", "Bilinmiyor")
    ozet = str(row.get("Odaklı Özet", "Özet yok"))

    if title:
        title.text = f"Sayfa {page_info} - {kategori}"
    if content:
        content.text = ozet # type: ignore

prs.save(OUTPUT_PPTX)
print("✅ PowerPoint sunumu oluşturuldu:", OUTPUT_PPTX)
