import os
import json
import fitz  # (PyMuPDF) PDF okuma
import matplotlib.pyplot as plt
import pandas as pd
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfbase import pdfmetrics
from pptx import Presentation
from pptx.util import Inches, Pt
from langchain_ollama import ChatOllama

# =====================
# Dosya yolları
# =====================
pdf_path = "data/Emsal-Karar.pdf"
parsed_file = "parsed.json"
summarized_file = "summarized.json"
stats_file = "emsal_stats.json"

# =====================
# PDF Ayrıştırma
# =====================
def parse_pdf():
    doc = fitz.open(pdf_path)
    pages = []
    for i, page in enumerate(doc, start=1): # type: ignore
        text = page.get_text("text")
        pages.append({"page": i, "text": text.strip()})
    with open(parsed_file, "w", encoding="utf-8") as f:
        json.dump(pages, f, ensure_ascii=False, indent=2)
    print(f"✅ PDF ayrıştırıldı, '{parsed_file}' dosyası oluşturuldu.")

# =====================
# Özetleme
# =====================
def summarize():
    with open(parsed_file, "r", encoding="utf-8") as f:
        pages = json.load(f)

    llm = ChatOllama(model="gemma3:4b")
    summaries = []
    for page in pages:
        prompt = f"Şu içeriği özetle:\n\n{page['text']}"
        response = llm.invoke(prompt)
        summaries.append({
            "page": page["page"],
            "summary": response.content.strip() # type: ignore
        })

    with open(summarized_file, "w", encoding="utf-8") as f:
        json.dump(summaries, f, ensure_ascii=False, indent=2)
    print(f"✅ Özetleme tamamlandı: {summarized_file}")

def load_or_summarize():
    if not os.path.exists(summarized_file):
        print("ℹ️ İlk kez çalıştırılıyor, özetleme yapılıyor...")
        if not os.path.exists(parsed_file):
            parse_pdf()
        summarize()

# =====================
# Chatbot
# =====================
def chatbot():
    load_or_summarize()
    with open(summarized_file, "r", encoding="utf-8") as f:
        summaries = json.load(f)

    llm = ChatOllama(model="gemma3:4b")
    print("📖 PDF Chatbot hazır! Çıkmak için 'q' yaz.\n")

    while True:
        question = input("Sorunuzu yazın: ")
        if question.lower() == "q":
            break
        context = "\n".join(
            f"- Sayfa {item['page']}: {item['summary']}" for item in summaries
        )
        prompt = f"{context}\n\nSoru: {question}\nCevap:"
        response = llm.invoke(prompt)
        print(f"\n🤖 Cevap: {response.content.strip()}\n") # type: ignore

# =====================
# İstatistikler
# =====================
def generate_stats():
    load_or_summarize()
    with open(summarized_file, "r", encoding="utf-8") as f:
        summaries = json.load(f)

    stats = {
        "Toplam Sayfa": len(summaries),
        "Özeti Olan Sayfa": sum(1 for s in summaries if s["summary"]),
        "Özetsiz Sayfa": sum(1 for s in summaries if not s["summary"]),
    }

    lengths = [len(s["summary"]) for s in summaries if s["summary"]]
    if lengths:
        stats["Ortalama Özet Uzunluğu"] = sum(lengths) / len(lengths) # type: ignore
        stats["En Kısa Özet Uzunluğu"] = min(lengths)
        stats["En Uzun Özet Uzunluğu"] = max(lengths)

    for s in summaries:
        stats[f"Sayfa {s['page']}"] = len(s["summary"])

    with open(stats_file, "w", encoding="utf-8") as f:
        json.dump(stats, f, ensure_ascii=False, indent=2)

    print("📊 PDF İstatistikleri:")
    for k, v in stats.items():
        print(f"- {k}: {v}")
    print(f"\n✅ İstatistikler kaydedildi: {stats_file}")

# =====================
# Grafik
# =====================
def generate_graph():
    if not os.path.exists(stats_file):
        generate_stats()

    with open(stats_file, "r", encoding="utf-8") as f:
        stats = json.load(f)

    pages = [k for k in stats.keys() if k.startswith("Sayfa")]
    lengths = [stats[k] for k in pages]

    plt.figure(figsize=(6, 4))
    plt.bar(pages, lengths)
    plt.title("Sayfa Bazlı Özet Uzunlukları")
    plt.xlabel("Sayfa")
    plt.ylabel("Özet Uzunluğu (karakter)")
    plt.xticks(rotation=45)
    plt.tight_layout()
    plt.savefig("emsal_summary_lengths.png")
    plt.close()
    print("✅ Grafik oluşturuldu: emsal_summary_lengths.png")

# =====================
# Rapor (PDF)
# =====================
def generate_report():
    if not os.path.exists(stats_file):
        generate_stats()

    with open(stats_file, "r", encoding="utf-8") as f:
        stats = json.load(f)

    pdfmetrics.registerFont(UnicodeCIDFont("HeiseiMin-W3"))
    doc = SimpleDocTemplate("emsal_report.pdf")
    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph("📑 Emsal Karar İstatistik Raporu", styles["Title"]))
    story.append(Spacer(1, 12))

    for k, v in stats.items():
        story.append(Paragraph(f"<b>{k}:</b> {v}", styles["Normal"]))
        story.append(Spacer(1, 6))

    doc.build(story)
    print("✅ Rapor oluşturuldu: emsal_report.pdf")

# =====================
# Sunum (PPTX)
# =====================
def generate_presentation():
    load_or_summarize()
    with open(summarized_file, "r", encoding="utf-8") as f:
        summaries = json.load(f)

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Emsal Karar Sunumu" # type: ignore
    slide.placeholders[1].text = "PDF özetlerinden oluşturulmuş sunum" # type: ignore

    for item in summaries:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Sayfa {item['page']}" # type: ignore
        slide.placeholders[1].text = item["summary"] # type: ignore

    prs.save("emsal_final_presentation.pptx")
    print("✅ Final sunum oluşturuldu: emsal_final_presentation.pptx")

# =====================
# Menü
# =====================
def menu():
    while True:
        print("\n=== PDFBot Final ===")
        print("1 - Chatbot (soru-cevap)")
        print("2 - İstatistikleri gör")
        print("3 - Grafik oluştur")
        print("4 - Rapor oluştur (PDF)")
        print("5 - Sunum oluştur (PPTX)")
        print("q - Çıkış")

        choice = input("Seçiminiz: ").strip()
        if choice == "1":
            chatbot()
        elif choice == "2":
            generate_stats()
        elif choice == "3":
            generate_graph()
        elif choice == "4":
            generate_report()
        elif choice == "5":
            generate_presentation()
        elif choice.lower() == "q":
            print("👋 Çıkılıyor...")
            break
        else:
            print("⚠️ Geçersiz seçim, tekrar deneyin.")

# =====================
# Ana Çalıştırma
# =====================
if __name__ == "__main__":
    menu()
