# pdfbot_17.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Sunum başlat
prs = Presentation()

# ---------- 1. Kapak ----------
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Emsal Karar Projesi - Final Sunumu" # type: ignore
subtitle.text = "4 Haftalık Çalışma Özeti\nHazırlayan: Esra" # type: ignore

# ---------- 2. Projenin Amacı ----------
slide_layout = prs.slide_layouts[1]  # Title + Content
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Projenin Amacı" # type: ignore
content = slide.placeholders[1].text_frame # type: ignore
content.text = "PDF tabanlı hukuk dokümanlarını işleyip sınıflandırmak, özetlemek ve etkileşimli chatbot ile kullanıcıya sunmak."

# ---------- 3. 1. Hafta ----------
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "1. Hafta: PDF Okuma ve Yapı Çözümleme" # type: ignore
content = slide.placeholders[1].text_frame # type: ignore
content.text = (
    "- PyMuPDF ile PDF okuma\n"
    "- Başlık, paragraf ayrımı\n"
    "- JSON formatında saklama\n"
    "- Test PDF ile deneme"
)

# ---------- 4. 2. Hafta ----------
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "2. Hafta: İçerik Sınıflandırma ve Etiketleme" # type: ignore
content = slide.placeholders[1].text_frame # type: ignore
content.text = (
    "- LLM ile kategori belirleme\n"
    "- Metinlere özet çıkarımı\n"
    "- JSON'a etiket ekleme\n"
    "- LangGraph classification node"
)

# ---------- 5. 3. Hafta ----------
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "3. Hafta: Chatbot Entegrasyonu" # type: ignore
content = slide.placeholders[1].text_frame # type: ignore
content.text = (
    "- LangGraph ile soru-cevap akışı\n"
    "- PDF'e göre anlık chat\n"
    "- Örnek sorularla test\n"
    "- Sayfa bazlı sorgulama"
)

# ---------- 6. 4. Hafta ----------
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "4. Hafta: Raporlama ve Sunum" # type: ignore
content = slide.placeholders[1].text_frame # type: ignore
content.text = (
    "- İstatistiklerin çıkarılması\n"
    "- Grafiksel analiz\n"
    "- PDF raporu\n"
    "- Final sunum hazırlanması"
)

# ---------- 7. Grafik Ekle ----------
slide_layout = prs.slide_layouts[5]  # Title Only
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Sayfa Bazlı Özet Uzunlukları" # type: ignore
slide.shapes.add_picture("emsal_summary_lengths.png", Inches(2), Inches(1.5), width=Inches(6))

# ---------- 8. Sonuç ----------
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Sonuç" # type: ignore
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = (
    "Bu proje kapsamında PDF’lerden anlamlı bilgi çıkarma, "
    "özetleme, chatbot entegrasyonu ve raporlama adımları başarıyla tamamlanmıştır."
)
p.font.size = Pt(20)
p.alignment = PP_ALIGN.CENTER

# Sunumu kaydet
prs.save("emsal_final_presentation.pptx")
print("✅ Final sunum oluşturuldu: emsal_final_presentation.pptx")
