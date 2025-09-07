# pdfbot_9.py
# 2. Hafta - 5. Gün
# Güvenli JSON'dan (emsal_safe.json) çıktı üretme

import json
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfbase import pdfmetrics

INPUT_JSON = "emsal_safe.json"
OUTPUT_CSV = "emsal_safe.csv"
OUTPUT_MD = "emsal_safe.md"
OUTPUT_PPTX = "emsal_safe.pptx"
OUTPUT_PDF = "emsal_safe.pdf"

# --------------------------
# 1. JSON'u oku
# --------------------------
with open(INPUT_JSON, "r", encoding="utf-8") as f:
    data = json.load(f)

df = pd.DataFrame(data)

# --------------------------
# 2. CSV ve Markdown'a kaydet
# --------------------------
df.to_csv(OUTPUT_CSV, index=False, encoding="utf-8-sig")
with open(OUTPUT_MD, "w", encoding="utf-8") as f:
    f.write(df.to_markdown(index=False))

# --------------------------
# 3. PowerPoint üret
# --------------------------
prs = Presentation()

for _, row in df.iterrows():
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Başlık + içerik
    title = slide.shapes.title
    content = slide.placeholders[1]

    page_info = row.get("page", "Genel")
    label = row.get("label", "Bilinmiyor")
    summary = row.get("summary", "Özet yok")

    title.text = f"Sayfa {page_info} - {label}" # type: ignore
    content.text = str(summary) # type: ignore

prs.save(OUTPUT_PPTX)

# --------------------------
# 4. PDF üret
# --------------------------
pdfmetrics.registerFont(UnicodeCIDFont("HeiseiKakuGo-W5"))  # Türkçe için destekli font
doc = SimpleDocTemplate(OUTPUT_PDF)
styles = getSampleStyleSheet()
story = []

story.append(Paragraph("Emsal Karar Güvenli Raporu", styles["Title"]))
story.append(Spacer(1, 20))

for _, row in df.iterrows():
    story.append(Paragraph(f"<b>Sayfa {row.get('page', 'Genel')} - {row.get('label', 'Bilinmiyor')}</b>", styles["Heading2"]))
    story.append(Paragraph(str(row.get("summary", "Özet yok")), styles["Normal"]))
    story.append(Spacer(1, 12))

doc.build(story)

print("✅ Güvenli çıktı üretildi:")
print("-", OUTPUT_CSV)
print("-", OUTPUT_MD)
print("-", OUTPUT_PPTX)
print("-", OUTPUT_PDF)
