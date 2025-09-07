# pdfbot_12.py
# Proje özeti sunumu oluşturucu
# - Mevcut dosyaları kontrol eder
# - Her hafta için özet slaytları oluşturur
# - Örnek Q&A slaytları ekler

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
import pandas as pd
import datetime
import os

# ---------- Ayarlar / Dosya yolları ----------
OUT_PPTX = "emsal_project_summary.pptx"
SAFE_JSON = Path("emsal_safe.json")
FOCUSED_CSV = Path("emsal_focused_table.csv")
REPORT_MD = Path("emsal_report.md")
REPORT_PDF = Path("emsal_report.pdf")

# ---------- Yardımcı fonksiyonlar ----------
def add_title_slide(prs: Presentation, title: str, subtitle: str): # type: ignore
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    if len(slide.placeholders) > 1 and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle

def add_bullet_slide(prs: Presentation, title: str, bullets): # type: ignore
    layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    # İlk satır
    if bullets:
        tf.text = bullets[0]
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
            p.level = 1

def add_qna_slide(prs: Presentation, q: str, a: str): # type: ignore
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Örnek Soru-Cevap"
    tf = slide.placeholders[1].text_frame
    tf.text = f"Soru: {q}\n\nCevap: {a}"

def short(s: str, n=400):
    if not s:
        return ""
    s2 = " ".join(s.split())
    return s2 if len(s2) <= n else s2[:n].rstrip() + "..."

# ---------- Sunum oluştur ----------
prs = Presentation()

# Başlık slaytı
project_title = "PDFBot Proje Özeti — Haftalar 1–3"
subtitle = f"Oluşturma tarihi: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}"
add_title_slide(prs, project_title, subtitle)

# Hafta 1 özeti
week1 = [
    "Amaç: PDF içeriğini okumak ve anlamlı parçalara ayırmak.",
    "Araçlar: PyMuPDF (fitz) ile sayfa/paragraf çıkartma.",
    "Çıktılar: emsal_parsed.json (ham metin), temel paragraf/sayfa ayrımı.",
    "Not: LangGraph başlangıcı 3. haftaya bırakıldı."
]
add_bullet_slide(prs, "1. Hafta — PDF Okuma ve Yapı Çözümleme", week1)

# Hafta 1 çıktı dosyaları
files_week1 = []
if REPORT_PDF.exists():
    files_week1.append(f"- PDF rapor: {REPORT_PDF.name}")
if SAFE_JSON.exists():
    files_week1.append(f"- Güvenli JSON (orem): {SAFE_JSON.name}")
if Path("emsal_parsed.json").exists():
    files_week1.append("- emsal_parsed.json")
if not files_week1:
    files_week1 = ["(Bu klasörde Hafta 1 çıktısı bulunamadı)"]
add_bullet_slide(prs, "1. Hafta — Oluşan Dosyalar", files_week1)

# Hafta 2 özeti
week2 = [
    "Amaç: LLM ile içerik sınıflandırma ve özetleme.",
    "Yöntem: Ollama (gemma3:4b) ile blok bazlı sınıflandırma ve odaklı özetleme.",
    "Çıktılar: emsal_classified.json, emsal_focused.json, emsal_focused_table.csv",
    "Güvenlik: Zararlı içerikler temizlendi -> emsal_safe.json"
]
add_bullet_slide(prs, "2. Hafta — LLM ile Sınıflandırma ve Etiketleme", week2)

# Hafta 2 çıktı dosyaları
files_week2 = []
if Path("emsal_classified.json").exists():
    files_week2.append("- emsal_classified.json")
if Path("emsal_focused.json").exists():
    files_week2.append("- emsal_focused.json")
if FOCUSED_CSV.exists():
    files_week2.append(f"- {FOCUSED_CSV.name}")
if Path("emsal_report.md").exists():
    files_week2.append("- emsal_report.md")
if not files_week2:
    files_week2 = ["(Hafta 2 çıktısı bulunamadı)"]
add_bullet_slide(prs, "2. Hafta — Oluşan Dosyalar", files_week2)

# Hafta 3 özeti
week3 = [
    "Amaç: Chatbot entegrasyonu — kullanıcı PDF'e göre soru sorabilsin.",
    "Araçlar: LangGraph + langchain-ollama (gemma3:4b).",
    "Yapılanlar: Basit retrieval, sayfa-odaklı cevaplama, Q&A testleri.",
    "Çıktı: emsal_presentation.pptx (otomatik sunum), etkileşimli CLI chatbot"
]
add_bullet_slide(prs, "3. Hafta — Chatbot Entegrasyonu ve Soru-Cevap", week3)

# Hafta 3 çıktı dosyaları
files_week3 = []
if Path("emsal_presentation.pptx").exists():
    files_week3.append("- emsal_presentation.pptx")
files_week3.append(f"- {OUT_PPTX} (özeti bu dosyada)")
add_bullet_slide(prs, "3. Hafta — Oluşan Dosyalar", files_week3)

# Örnek Q&A slaytları (kullanıcı testlerinden alıntılar)
# Eğer csv varsa, örnekleri oradan al, yoksa sabit örnek kullan
example_qnas = []
if FOCUSED_CSV.exists():
    try:
        df = pd.read_csv(str(FOCUSED_CSV))
        # alttan 3 satırdan örnekler al
        for i, row in df.head(5).iterrows():
            q = f"Sayfa {row.get('Sayfa', 'Genel')} hakkında"
            a = short(str(row.get('Odaklı Özet', 'Özet yok')), 300)
            example_qnas.append((q, a))
    except Exception:
        example_qnas = []
if not example_qnas:
    # Kullanıcının testlerinden alınmış örnek cevaplar (sık kullanılan sorular)
    example_qnas = [
        ("Bu dava neyle ilgilidir?",
         "Dava, arabuluculuk tutanağının icra edilebilirliğinin değerlendirilmesi ve icra edilebilirlik erhi talebinin reddi ile ilgilidir."),
        ("Sayfa 2'de karar neydi?",
         "Yerel mahkeme, arabuluculuk tutanağının icra edilebilirlik şartlarını taşımadığı gerekçesiyle talebi reddetmiştir."),
        ("Davacı kimdi ve ne talep etti?",
         "Davacı: AT VE TİC LTD ŞTİ.; Talep: arabuluculuk tutanağına icra edilebilirlik şerhi verilmesi.")
    ]

# Her Q&A için bir slayt
for q,a in example_qnas:
    add_qna_slide(prs, q, a)

# Mimari / Akış slaytı
arch = [
    "1) PDF (Emsal-Karar.pdf) → PyMuPDF ile metin çıkarma",
    "2) Metin → bölümlere ayırma → JSON (emsal_parsed.json)",
    "3) LLM sınıflandırma (pdfbot_2.py) → emsal_classified.json",
    "4) Odaklı özetleme (pdfbot_5.py) → emsal_focused.json",
    "5) Güvenlik filtresi (pdfbot_8.py) → emsal_safe.json",
    "6) Chatbot (pdfbot_11.py) : LangGraph + Gemma3 ile Soru-Cevap",
    "7) Sunum/rapor üretimi (pdfbot_3/4/9/12)"
]
add_bullet_slide(prs, "Proje Mimari Akışı", arch)

# Nasıl çalıştırılır (komutlar)
commands = [
    "1) Sanal ortam aktif edin: .\\venv\\Scripts\\Activate.ps1",
    "2) Gerekli paketleri kurun (örnek): python -m pip install python-pptx pandas langchain-ollama langgraph",
    "3) Sırasıyla çalıştır:",
    "   python pdfbot_1.py  # parse",
    "   python pdfbot_2.py  # sınıflandırma",
    "   python pdfbot_3.py  # markdown/csv",
    "   python pdfbot_4.py  # pdf rapor",
    "   python pdfbot_5.py  # odaklı özetleme",
    "   python pdfbot_6.py  # tablo",
    "   python pdfbot_7.py  # sunum",
    "   python pdfbot_8.py  # güvenlik filtre",
    "   python pdfbot_9.py  # güvenli çıktı üret",
    "   python pdfbot_11.py # chatbot (etkileşimli)"
]
add_bullet_slide(prs, "Çalıştırma - Hızlı Komutlar", commands)

# Troubleshooting
troubles = [
    "pip hatası: pip yerine 'python -m pip install ...' kullanın (venv yol hatası için).",
    "Model bulunamadı: 'ollama list' ile model adını kontrol edin (ör: gemma3:4b).",
    "Türkçe karakter hatası: reportlab için sistem fontu kullanın (Arial veya DejaVu) ve pdfbot_4.py'yi güncelleyin.",
    "Pandas tabulate hatası: 'python -m pip install tabulate' ile giderilir.",
    "VSCode uyarıları: Python interpreter'ı venv'e ayarlayın (Python: Select Interpreter)."
]
add_bullet_slide(prs, "Common Troubleshooting (Yaygın Sorunlar)", troubles)

# Next steps / gelişim maddeleri
next_steps = [
    "1) Chat arayüzü: web veya masaüstü GUI (Streamlit / Flask).",
    "2) Daha sofistike retriever: BM25 veya embedding tabanlı arama (faiss).",
    "3) Slaytlara otomatik tablo/diagram ekleme (matplotlib + resim ekleme).",
    "4) Kullanıcı kimlik doğrulama ve veri gizliliği iyileştirmeleri."
]
add_bullet_slide(prs, "Next Steps (Geliştirme Önerileri)", next_steps)

# Kapanış slaytı
add_bullet_slide(prs, "Kapanış ve Sorular", ["Sunum hazır — aklınıza takılan her şeyi sorun, adım adım açıklarım."])

# Kaydet
prs.save(OUT_PPTX)
print("✅ Sunum oluşturuldu:", OUT_PPTX)
